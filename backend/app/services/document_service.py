from pathlib import Path
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


def extract_text(file_path: str) -> tuple[str, int]:
    """Extract text from PDF, DOCX, or PPTX. Returns (text, page_count)."""
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def extract_pdf(file_path: str) -> tuple[str, int]:
    """Extract text from PDF using PyMuPDF."""
    doc = fitz.open(file_path)
    text_parts = []
    page_count = doc.page_count
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n\n".join(text_parts), page_count


def extract_docx(file_path: str) -> tuple[str, int]:
    """Extract text from DOCX."""
    doc = Document(file_path)
    text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
    page_count = len(doc.sections) or 1
    return "\n\n".join(text_parts), page_count


def extract_pptx(file_path: str) -> tuple[str, int]:
    """Extract text from PPTX."""
    prs = Presentation(file_path)
    text_parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
    return "\n\n".join(text_parts), len(prs.slides)


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> list[str]:
    """Split text into overlapping chunks of approximately chunk_size tokens."""
    # Simple word-based chunking (approximation: 1 token ~= 0.75 words)
    words = text.split()
    word_chunk_size = int(chunk_size * 0.75)
    word_overlap = int(overlap * 0.75)

    chunks = []
    start = 0
    while start < len(words):
        end = start + word_chunk_size
        chunk = " ".join(words[start:end])
        if chunk.strip():
            chunks.append(chunk)
        start += word_chunk_size - word_overlap
    return chunks
